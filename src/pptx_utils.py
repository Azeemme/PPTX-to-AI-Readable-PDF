"""Extract speaker notes, slide count, and alt-text from PPTX using python-pptx."""

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.slide import Slide


def _notes_text(slide: Slide) -> str:
    """Get speaker notes text for one slide; empty string if none."""
    if not slide.has_notes_slide:
        return ""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        return (tf.text or "").strip()
    except Exception:
        return ""


def get_speaker_notes(pptx_path: str | Path) -> list[str]:
    """
    Return a list of speaker note strings, one per slide, in order.
    Empty string for slides with no notes.
    """
    path = Path(pptx_path)
    if not path.is_file():
        return []
    try:
        prs = Presentation(str(path))
        return [_notes_text(s) for s in prs.slides]
    except Exception:
        return []


def get_slide_count(pptx_path: str | Path) -> int:
    """Return the number of slides in the presentation."""
    path = Path(pptx_path)
    if not path.is_file():
        return 0
    try:
        prs = Presentation(str(path))
        return len(prs.slides)
    except Exception:
        return 0


def get_slide_text_and_alt(pptx_path: str | Path) -> list[dict[str, Any]]:
    """
    Return per-slide text and alt-text for embedding in context.
    Each item: {"text": str, "alt": str}.
    """
    path = Path(pptx_path)
    if not path.is_file():
        return []
    result: list[dict[str, Any]] = []
    try:
        prs = Presentation(str(path))
        for slide in prs.slides:
            texts: list[str] = []
            alt_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
                if hasattr(shape, "alt_text") and shape.alt_text:
                    alt_texts.append(shape.alt_text.strip())
            result.append({
                "text": "\n\n".join(filter(None, texts)),
                "alt": " ".join(filter(None, alt_texts)),
            })
        return result
    except Exception:
        return []
